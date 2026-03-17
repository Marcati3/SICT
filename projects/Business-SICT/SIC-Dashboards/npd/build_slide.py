#!/usr/bin/env python3
"""
Build executive board slide: New Product Release Health Summary
Data source: New_Release_Product_Health_Dashboard.xlsx
Output: Single PowerPoint slide for C-suite / Board presentation
"""

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
XLSX = os.path.join(SCRIPT_DIR, "New_Release_Product_Health_Dashboard.xlsx")
OUTPUT = os.path.join(SCRIPT_DIR, "NPD_Board_Slide_Mar2026.pptx")

# ── Colors ──
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
BORDER = RGBColor(0x33, 0x41, 0x55)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)

RAG_COLORS = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED, "N/A": MUTED}
RAG_LABELS = {"GREEN": "On Track", "AMBER": "At Risk", "RED": "Action Req'd", "N/A": "N/A"}


def load_data():
    wb = openpyxl.load_workbook(XLSX, data_only=True)

    # Executive Summary sheet — product rows start at row 8 (row 7 is header)
    ws = wb["Executive Summary"]
    products = []
    def num(val, default=0):
        if val is None or val == "N/A" or val == "":
            return default
        try:
            return float(val)
        except (ValueError, TypeError):
            return default

    for row in ws.iter_rows(min_row=8, max_row=17, values_only=True):
        if row[0] is None:
            continue
        # Skip header rows
        if row[0] == "Project Name" or row[0] == "Product":
            continue
        products.append({
            "name": row[0],
            "chip": row[1],
            "launch": row[2],
            "gate": row[3],
            "status": row[4],
            "investment": num(row[5]),
            "bc_rev": num(row[6]),
            "actual_rev": num(row[7]),
            "variance": num(row[8]),
            "trajectory_pct": num(row[9]),
            "gp_pct": num(row[10]),
            "bc_payback": row[11],
            "rag": row[12] or "N/A",
            "whats_off": row[13] or "",
            "why": row[14] or "",
            "action": row[15] or "",
            "gp_rag": row[16] or "GREEN",
            "payback_signal": row[17] or "",
        })

    # LEAN Scorecard
    ws_lean = wb["LEAN Scorecard"]
    lean = []
    for row in ws_lean.iter_rows(min_row=4, max_row=13, values_only=True):
        if row[0] is None:
            continue
        lean.append({
            "name": row[0],
            "rev_traj": row[1] or "N/A",
            "gp_vs_target": row[2] or "N/A",
            "payback": row[3] or "N/A",
            "overall_rag": row[4] or "N/A",
            "whats_off": row[5] or "",
            "why": row[6] or "",
            "action": row[7] or "",
            "total_score": row[11] or 0,
        })

    # Forward forecast from Revenue Tracking
    ws_rev = wb["Revenue Tracking"]
    fwd = []
    for row in ws_rev.iter_rows(min_row=14, max_row=23, values_only=True):
        if row[0] is None or row[0] == "FORWARD FORECAST (Biz Case Projections 2026-2030)":
            continue
        if row[0] == "Product":
            continue
        fwd.append({
            "name": row[0],
            "chip": row[1],
            "fwd_total": row[7] or 0,
            "cumulative": row[8] or 0,
            "investment": row[9] or 0,
            "proj_roi": row[10] or 0,
        })

    wb.close()
    return products, lean, fwd


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    # Subtle rounding
    shape.adjustments[0] = 0.05
    return shape


def add_text(slide, left, top, width, height, text, font_size=10,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def fmt_thb(val):
    if val is None or val == "N/A":
        return "N/A"
    v = float(val)
    if abs(v) >= 1_000_000:
        return f"฿{v/1_000_000:.1f}M"
    elif abs(v) >= 1_000:
        return f"฿{v/1_000:.0f}K"
    return f"฿{v:.0f}"


def fmt_pct(val):
    if val is None or val == "N/A" or val == 0:
        return "N/A"
    return f"{float(val)*100:.0f}%"


def build_slide(products, lean, fwd):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # ── Background ──
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG

    # ── HEADER BAR ──
    add_rounded_rect(slide, Inches(0), Inches(0), Inches(16), Inches(0.85), RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, Inches(0.4), Inches(0.12), Inches(8), Inches(0.4),
             "New Product Release — Health Summary", 20, WHITE, True)
    add_text(slide, Inches(0.4), Inches(0.48), Inches(10), Inches(0.3),
             "LEAN Management Assessment  |  RAG = Trajectory vs Plan  |  Currency: THB  |  FX: 33 THB/USD  |  March 2026",
             9, MUTED)

    # Header right — portfolio counts
    active = sum(1 for p in products if p["status"] == "Active")
    research = sum(1 for p in products if p["status"] == "Research")
    total_inv = sum(p["investment"] for p in products)
    total_actual = sum(p["actual_rev"] for p in products)

    rag_counts = {"GREEN": 0, "AMBER": 0, "RED": 0}
    for p in products:
        if p["rag"] in rag_counts:
            rag_counts[p["rag"]] += 1

    add_text(slide, Inches(10.5), Inches(0.12), Inches(5.2), Inches(0.35),
             f"{len(products)} Products  |  {active} Active  |  {research} Research  |  "
             f"Investment: {fmt_thb(total_inv)}  |  Actual Rev: {fmt_thb(total_actual)}",
             10, LIGHT, False, PP_ALIGN.RIGHT)

    add_text(slide, Inches(10.5), Inches(0.48), Inches(5.2), Inches(0.3),
             f"🟢 {rag_counts['GREEN']} On Track    🟡 {rag_counts['AMBER']} At Risk    🔴 {rag_counts['RED']} Action Req'd",
             9, LIGHT, False, PP_ALIGN.RIGHT)

    # ═══════════════════════════════════════════════════════
    # ROW 1: KPI SCORECARDS (4 cards across top)
    # ═══════════════════════════════════════════════════════
    y_kpi = Inches(1.05)
    card_w = Inches(3.7)
    card_h = Inches(1.15)
    gap = Inches(0.2)

    # Card 1: Total Investment
    x = Inches(0.4)
    add_rounded_rect(slide, x, y_kpi, card_w, card_h, CARD_BG, BORDER)
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.08), card_w, Inches(0.2),
             "TOTAL PORTFOLIO INVESTMENT", 8, MUTED, True)
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.32), card_w, Inches(0.35),
             fmt_thb(total_inv), 22, ACCENT_BLUE, True)
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.72), card_w, Inches(0.2),
             f"${total_inv/33/1_000_000:.1f}M USD  |  {len(products)} products", 9, LIGHT)

    # Card 2: Actual Revenue (2023-2025)
    x = Inches(0.4) + card_w + gap
    add_rounded_rect(slide, x, y_kpi, card_w, card_h, CARD_BG, BORDER)
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.08), card_w, Inches(0.2),
             "ACTUAL REVENUE (2023–2025)", 8, MUTED, True)
    rev_color = GREEN if total_actual > 0 else MUTED
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.32), card_w, Inches(0.35),
             fmt_thb(total_actual), 22, rev_color, True)
    total_bc = sum(p["bc_rev"] for p in products)
    variance = total_actual - total_bc
    var_color = GREEN if variance >= 0 else RED
    var_sign = "+" if variance >= 0 else ""
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.72), card_w, Inches(0.2),
             f"vs BC {fmt_thb(total_bc)}  |  Var: {var_sign}{fmt_thb(variance)}", 9, LIGHT)

    # Card 3: Forward Revenue (2026-2030)
    x = Inches(0.4) + 2 * (card_w + gap)
    total_fwd = sum(f["fwd_total"] for f in fwd)
    total_cum = sum(f["cumulative"] for f in fwd)
    add_rounded_rect(slide, x, y_kpi, card_w, card_h, CARD_BG, BORDER)
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.08), card_w, Inches(0.2),
             "FORWARD REVENUE (2026–2030 BC)", 8, MUTED, True)
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.32), card_w, Inches(0.35),
             fmt_thb(total_fwd), 22, GREEN, True)
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.72), card_w, Inches(0.2),
             f"Lifetime cumulative: {fmt_thb(total_cum)}", 9, LIGHT)

    # Card 4: Portfolio GP (weighted avg)
    x = Inches(0.4) + 3 * (card_w + gap)
    active_prods = [p for p in products if p["status"] == "Active" and p["actual_rev"] > 0]
    if active_prods:
        weighted_gp_num = sum(p["gp_pct"] * p["actual_rev"] for p in active_prods)
        weighted_gp_den = sum(p["actual_rev"] for p in active_prods)
        avg_gp = weighted_gp_num / weighted_gp_den if weighted_gp_den else 0
    else:
        avg_gp = 0
    add_rounded_rect(slide, x, y_kpi, card_w, card_h, CARD_BG, BORDER)
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.08), card_w, Inches(0.2),
             "WEIGHTED AVG GROSS PROFIT %", 8, MUTED, True)
    gp_color = GREEN if avg_gp >= 0.50 else AMBER if avg_gp >= 0.35 else RED
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.32), card_w, Inches(0.35),
             f"{avg_gp*100:.1f}%", 22, gp_color, True)
    gp_ambers = sum(1 for p in products if p["gp_rag"] == "AMBER")
    gp_reds = sum(1 for p in products if p["gp_rag"] == "RED")
    add_text(slide, x + Inches(0.15), y_kpi + Inches(0.72), card_w, Inches(0.2),
             f"Revenue-weighted across active products  |  {gp_ambers} Amber, {gp_reds} Red", 9, LIGHT)

    # ═══════════════════════════════════════════════════════
    # ROW 2: PRODUCT PORTFOLIO TABLE (main table)
    # ═══════════════════════════════════════════════════════
    y_table = Inches(2.4)
    add_text(slide, Inches(0.4), y_table - Inches(0.28), Inches(4), Inches(0.25),
             "01  PRODUCT PORTFOLIO — LEAN SCORECARD", 9, MUTED, True)

    # Table: 10 products x key columns
    active_products = [p for p in products if p["status"] == "Active"]
    research_products = [p for p in products if p["status"] == "Research"]
    all_sorted = active_products + research_products

    # Find matching lean data
    lean_map = {l["name"]: l for l in lean}

    cols = ["Product", "Chip", "Gate", "Investment", "BC Rev", "Actual Rev",
            "Traj %", "GP %", "Overall", "Score"]
    col_widths = [Inches(1.75), Inches(0.75), Inches(0.7), Inches(1.0), Inches(1.0),
                  Inches(1.0), Inches(0.65), Inches(0.6), Inches(0.7), Inches(0.55)]
    total_table_w = sum(cw for cw in col_widths)

    tbl_shape = slide.shapes.add_table(
        len(all_sorted) + 1, len(cols),
        Inches(0.4), y_table, total_table_w, Inches(0.25) * (len(all_sorted) + 1)
    )
    tbl = tbl_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    # Style header row
    for j, col_name in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(7.5)
            paragraph.font.color.rgb = MUTED
            paragraph.font.bold = True
            paragraph.font.name = "Segoe UI"
            paragraph.alignment = PP_ALIGN.LEFT

    # Data rows
    for i, p in enumerate(all_sorted):
        row_idx = i + 1
        lean_data = lean_map.get(p["name"], {})
        overall_rag = lean_data.get("overall_rag", p["rag"])
        score = lean_data.get("total_score", 0)

        values = [
            p["name"],
            p["chip"],
            p["gate"],
            fmt_thb(p["investment"]),
            fmt_thb(p["bc_rev"]) if p["bc_rev"] else "—",
            fmt_thb(p["actual_rev"]) if p["actual_rev"] else "—",
            fmt_pct(p["trajectory_pct"]),
            fmt_pct(p["gp_pct"]),
            overall_rag,
            f"{score:.0f}" if score else "—",
        ]

        for j, val in enumerate(values):
            cell = tbl.cell(row_idx, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else RGBColor(0x17, 0x22, 0x34)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(7.5)
                paragraph.font.name = "Segoe UI"
                paragraph.font.color.rgb = LIGHT
                paragraph.alignment = PP_ALIGN.LEFT
                # Color the RAG column
                if j == 8:  # Overall RAG
                    paragraph.font.color.rgb = RAG_COLORS.get(val, MUTED)
                    paragraph.font.bold = True
                # Color trajectory
                if j == 6 and p["trajectory_pct"]:
                    t = float(p["trajectory_pct"])
                    if t >= 0.8:
                        paragraph.font.color.rgb = GREEN
                    elif t >= 0.5:
                        paragraph.font.color.rgb = AMBER
                    elif t > 0:
                        paragraph.font.color.rgb = RED

    # ═══════════════════════════════════════════════════════
    # ROW 3 LEFT: ACTION ITEMS (Amber/Red commentary)
    # ═══════════════════════════════════════════════════════
    y_actions = Inches(5.6)
    action_w = Inches(7.7)
    add_text(slide, Inches(0.4), y_actions - Inches(0.28), Inches(4), Inches(0.25),
             "02  ACTIONS REQUIRED — NO ORPHAN REDS", 9, MUTED, True)

    action_items = [p for p in products if p["rag"] in ("RED", "AMBER") and p["action"]]
    action_items.sort(key=lambda x: 0 if x["rag"] == "RED" else 1)

    # Add LEAN-level red items that might not be in exec summary
    lean_actions = [l for l in lean if l["overall_rag"] in ("RED", "AMBER") and l["action"]]
    seen_names = {a["name"] for a in action_items}
    for l in lean_actions:
        if l["name"] not in seen_names:
            action_items.append({
                "name": l["name"], "rag": l["overall_rag"],
                "whats_off": l["whats_off"], "why": l["why"], "action": l["action"]
            })

    if action_items:
        action_tbl_shape = slide.shapes.add_table(
            len(action_items) + 1, 4,
            Inches(0.4), y_actions,
            action_w, Inches(0.25) * (len(action_items) + 1)
        )
        action_tbl = action_tbl_shape.table
        action_col_widths = [Inches(0.55), Inches(1.5), Inches(2.5), Inches(3.15)]
        for j, w in enumerate(action_col_widths):
            action_tbl.columns[j].width = w

        for j, hdr in enumerate(["RAG", "Product", "What's Off", "Action (Owner / Date)"]):
            cell = action_tbl.cell(0, j)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(7.5)
                paragraph.font.color.rgb = MUTED
                paragraph.font.bold = True
                paragraph.font.name = "Segoe UI"

        for i, item in enumerate(action_items):
            row_idx = i + 1
            rag_str = item["rag"]
            values = [rag_str, item["name"], item.get("whats_off", ""), item.get("action", "")]
            for j, val in enumerate(values):
                cell = action_tbl.cell(row_idx, j)
                cell.text = str(val) if val else ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else RGBColor(0x17, 0x22, 0x34)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(7)
                    paragraph.font.name = "Segoe UI"
                    paragraph.font.color.rgb = LIGHT
                    if j == 0:
                        paragraph.font.color.rgb = RAG_COLORS.get(rag_str, MUTED)
                        paragraph.font.bold = True

    # ═══════════════════════════════════════════════════════
    # ROW 3 RIGHT: FORWARD REVENUE TOP 5 + INVESTMENT RECOVERY
    # ═══════════════════════════════════════════════════════
    fwd_x = Inches(8.3)
    fwd_w = Inches(7.3)
    add_text(slide, fwd_x, y_actions - Inches(0.28), Inches(4), Inches(0.25),
             "03  FORWARD OUTLOOK — BIZ CASE PROJECTIONS (2026–2030)", 9, MUTED, True)

    fwd_sorted = sorted(fwd, key=lambda f: f["fwd_total"], reverse=True)

    fwd_tbl_shape = slide.shapes.add_table(
        len(fwd_sorted) + 1, 5,
        fwd_x, y_actions,
        fwd_w, Inches(0.25) * (len(fwd_sorted) + 1)
    )
    fwd_tbl = fwd_tbl_shape.table
    fwd_col_widths = [Inches(1.75), Inches(0.75), Inches(1.3), Inches(1.3), Inches(2.2)]
    for j, w in enumerate(fwd_col_widths):
        fwd_tbl.columns[j].width = w

    for j, hdr in enumerate(["Product", "Chip", "Fwd Rev (BC)", "Cumulative", "Proj ROI"]):
        cell = fwd_tbl.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(7.5)
            paragraph.font.color.rgb = MUTED
            paragraph.font.bold = True
            paragraph.font.name = "Segoe UI"

    for i, f in enumerate(fwd_sorted):
        row_idx = i + 1
        roi_str = f"{f['proj_roi']:.1f}x" if f["proj_roi"] else "—"
        values = [f["name"], f["chip"], fmt_thb(f["fwd_total"]),
                  fmt_thb(f["cumulative"]), roi_str]
        for j, val in enumerate(values):
            cell = fwd_tbl.cell(row_idx, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else RGBColor(0x17, 0x22, 0x34)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(7)
                paragraph.font.name = "Segoe UI"
                paragraph.font.color.rgb = LIGHT
                if j == 4 and f["proj_roi"] and f["proj_roi"] > 5:
                    paragraph.font.color.rgb = GREEN
                    paragraph.font.bold = True

    # ═══════════════════════════════════════════════════════
    # FOOTER
    # ═══════════════════════════════════════════════════════
    add_text(slide, Inches(0.4), Inches(8.55), Inches(8), Inches(0.3),
             "Silicon Craft Technology PLC (SET: SICT)  |  Confidential — Board Use Only  |  Prepared: March 2026",
             7.5, MUTED, False, PP_ALIGN.LEFT)
    add_text(slide, Inches(8), Inches(8.55), Inches(7.6), Inches(0.3),
             "RAG = Trajectory vs Plan (not snapshot)  |  No Orphan Reds  |  LEAN Scoring: Rev 40% + GP 30% + ROI 30%",
             7.5, MUTED, False, PP_ALIGN.RIGHT)

    prs.save(OUTPUT)
    print(f"✓ Saved: {OUTPUT}")


if __name__ == "__main__":
    products, lean, fwd = load_data()
    build_slide(products, lean, fwd)
