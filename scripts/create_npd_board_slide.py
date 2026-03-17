"""
Generate 1-page NPD Health Dashboard slide for Board Executive Meeting.
Silicon Craft Technology (SICT) — March 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)       # Navy
HEADER_BG = RGBColor(0x0D, 0x47, 0xA1)      # Deep blue
CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)        # Light gray card
BORDER_GRAY = RGBColor(0xDE, 0xE2, 0xE6)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_MID = RGBColor(0x4A, 0x5A, 0x6A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xEF, 0x6C, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
AMBER_BG = RGBColor(0xFF, 0xF3, 0xE0)
RED_BG = RGBColor(0xFF, 0xEB, 0xEE)
ACCENT = RGBColor(0x15, 0x65, 0xC0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Helper functions ──
def add_rect(left, top, w, h, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.02
    return shape

def add_text(left, top, w, h, text, size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_kpi_card(left, top, w, h, label, value, sub="", rag_color=None):
    add_rect(left, top, w, h, fill=WHITE, border=BORDER_GRAY)
    # accent bar at top
    bar_color = rag_color if rag_color else ACCENT
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05), w - Inches(0.1), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    # label
    add_text(left + Inches(0.15), top + Inches(0.12), w - Inches(0.3), Inches(0.3),
             label, size=8, color=TEXT_MID, bold=False, align=PP_ALIGN.CENTER)
    # value
    add_text(left + Inches(0.15), top + Inches(0.35), w - Inches(0.3), Inches(0.4),
             value, size=20, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    # sub
    if sub:
        add_text(left + Inches(0.15), top + Inches(0.7), w - Inches(0.3), Inches(0.25),
                 sub, size=7, color=TEXT_MID, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# HEADER BAR
# ═══════════════════════════════════════════════════════════════
header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.75))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BG
header.line.fill.background()

add_text(Inches(0.5), Inches(0.1), Inches(8), Inches(0.35),
         "NEW PRODUCT RELEASE — HEALTH DASHBOARD", size=18, bold=True, color=WHITE)
add_text(Inches(0.5), Inches(0.42), Inches(8), Inches(0.25),
         "Silicon Craft Technology  |  Board of Directors  |  March 2026  |  LEAN Management Assessment  |  Currency: THB  |  FX: 33 THB/USD",
         size=8, color=RGBColor(0xB0, 0xBE, 0xC5))

# Confidential tag
add_text(Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.35),
         "CONFIDENTIAL", size=10, bold=True, color=RGBColor(0xFF, 0xAB, 0x40), align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════
# ROW 1: KPI SCORECARDS (y = 0.9)
# ═══════════════════════════════════════════════════════════════
y1 = Inches(0.9)
card_h = Inches(1.0)
card_w = Inches(2.0)
gap = Inches(0.15)
x_start = Inches(0.5)

add_kpi_card(x_start, y1, card_w, card_h,
             "TOTAL PRODUCTS", "10", "8 Active  |  2 Research", ACCENT)

add_kpi_card(x_start + (card_w + gap) * 1, y1, card_w, card_h,
             "TOTAL INVESTMENT", "114.7M", "THB (~$3.5M USD)", ACCENT)

add_kpi_card(x_start + (card_w + gap) * 2, y1, card_w, card_h,
             "ACTUAL REV (23-25)", "69.5M", "THB (~$2.1M USD)", ACCENT)

add_kpi_card(x_start + (card_w + gap) * 3, y1, card_w, card_h,
             "LIFETIME PORTFOLIO", "1.06B", "THB projected (~$32M USD)", GREEN)

add_kpi_card(x_start + (card_w + gap) * 4, y1, card_w, card_h,
             "PORTFOLIO GP (AVG)", "55%", "Range: 32% – 82%", ACCENT)

# RAG summary card (wider)
rag_x = x_start + (card_w + gap) * 5
rag_w = Inches(2.35)
add_rect(rag_x, y1, rag_w, card_h, fill=WHITE, border=BORDER_GRAY)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rag_x + Inches(0.05), y1 + Inches(0.05), rag_w - Inches(0.1), Inches(0.04))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()
add_text(rag_x + Inches(0.1), y1 + Inches(0.1), rag_w - Inches(0.2), Inches(0.2),
         "OVERALL RAG STATUS", size=8, color=TEXT_MID, align=PP_ALIGN.CENTER)

# RAG circles
def add_rag_dot(cx, cy, r, color, label, count):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(cx, cy + Inches(0.01), r, r - Inches(0.05),
             str(count), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(cx - Inches(0.05), cy + r + Inches(0.01), r + Inches(0.1), Inches(0.15),
             label, size=7, color=TEXT_MID, align=PP_ALIGN.CENTER)

dot_size = Inches(0.4)
dot_y = y1 + Inches(0.33)
dot_gap = Inches(0.7)
dot_x_start = rag_x + Inches(0.25)
add_rag_dot(dot_x_start, dot_y, dot_size, GREEN, "GREEN", 3)
add_rag_dot(dot_x_start + dot_gap, dot_y, dot_size, AMBER, "AMBER", 2)
add_rag_dot(dot_x_start + dot_gap * 2, dot_y, dot_size, RED, "RED", 3)


# ═══════════════════════════════════════════════════════════════
# ROW 2: PRODUCT TABLE (y = 2.05)
# ═══════════════════════════════════════════════════════════════
y2 = Inches(2.1)
tbl_w = Inches(12.333)
tbl_x = Inches(0.5)

# Section header
add_text(tbl_x, y2 - Inches(0.25), Inches(5), Inches(0.25),
         "PRODUCT PORTFOLIO — LEAN SCORECARD", size=10, bold=True, color=DARK_BG)

# Table data
cols = [1.7, 0.8, 0.7, 1.0, 0.9, 1.0, 0.9, 0.8, 0.7, 0.7, 3.2]  # total ~12.4
headers = ["Product", "Chip", "Launch", "Investment", "BC Rev", "Actual Rev", "Traj %", "GP %", "Rev RAG", "Overall", "Action Required"]
rows_data = [
    ["2201D-LFT-FOUP",    "SIC73F1",  "2023", "12.8M",  "73.4M",  "64.8M",  "88%",  "77%", "GREEN",  "AMBER", "Finance to review COGS vs BC (+9pp gap)"],
    ["2218D-SEN-SIC4343",  "SIC4343",  "2024", "1.5M",   "2.3M",   "64K",    "3%",   "82%", "RED",    "RED",   "Exec review: viability & market fit (immediate)"],
    ["2003D-LFT-SIC7805",  "SIC7150",  "2025", "29.4M",  "—",      "4.3M",   "100%", "33%", "GREEN",  "RED",   "Payback 4.5yr vs target 1.9yr — restructure review"],
    ["2406D-CNT-GT12",     "SIC279",   "2025", "4.8M",   "330K",   "0",      "0%",   "51%", "AMBER",  "AMBER", "Sales: customer update + revised ramp forecast"],
    ["2415D-CNT-Dragon",   "SIC72A1",  "2025", "1.2M",   "—",      "167K",   "100%", "55%", "GREEN",  "RED",   "Payback 4.9yr vs target 2.1yr — restructure review"],
    ["2417D-CNT-Marlin",   "SIC72A2",  "2025", "12.7M",  "—",      "164K",   "100%", "32%", "GREEN",  "GREEN", "Early stage — monitor"],
    ["2112D-CNT-NFCV",     "SIC56NL",  "2025", "34.9M",  "—",      "19K",    "100%", "49%", "GREEN",  "GREEN", "Early stage — monitor (largest investment)"],
    ["1805D-HDX3-SIC379",  "SIC379",   "2025", "17.1M",  "—",      "12K",    "100%", "68%", "GREEN",  "GREEN", "Early stage — monitor"],
]

n_cols = len(headers)
n_rows = len(rows_data) + 1
tbl_h = Inches(0.3 * n_rows)
table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_x, y2, tbl_w, tbl_h)
table = table_shape.table

# Col widths
for i, w in enumerate(cols):
    table.columns[i].width = Inches(w)

# Style header row
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BG
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(7.5)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# RAG color map
rag_map = {"GREEN": (GREEN, GREEN_BG), "AMBER": (AMBER, AMBER_BG), "RED": (RED, RED_BG)}

for r, row in enumerate(rows_data):
    for c, val in enumerate(row):
        cell = table.cell(r + 1, c)
        cell.text = val
        # Alternating rows
        if r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.color.rgb = TEXT_DARK
            p.alignment = PP_ALIGN.CENTER if c < 10 else PP_ALIGN.LEFT

        # Color RAG columns
        if c in (8, 9) and val in rag_map:
            fg, bg = rag_map[val]
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = fg
                p.font.bold = True

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)


# ═══════════════════════════════════════════════════════════════
# ROW 3: BOTTOM SECTION — KEY INSIGHTS + DECISIONS NEEDED (y = 4.9)
# ═══════════════════════════════════════════════════════════════
y3 = Inches(4.85)
box_h = Inches(2.4)

# ── Left box: Key Insights ──
left_w = Inches(5.9)
add_rect(tbl_x, y3, left_w, box_h, fill=WHITE, border=BORDER_GRAY)

add_text(tbl_x + Inches(0.2), y3 + Inches(0.1), Inches(3), Inches(0.22),
         "KEY INSIGHTS", size=10, bold=True, color=DARK_BG)

insights = [
    ("FOUP (SIC73F1) — Portfolio Anchor", "88% trajectory, THB 64.8M actual revenue, investment fully recovered (5.1x). GP 9pp below BC target — COGS review needed.", GREEN),
    ("SIC4343 — Critical Alert", "Revenue at 3% of plan (THB 2.2M gap). Structural miss, not timing. Requires immediate viability assessment.", RED),
    ("SIC7805 & Dragon — Payback Concern", "Revenue on track but payback pace 2-3x slower than BC target. High investment (THB 30.6M combined) needs recovery plan.", AMBER),
    ("6 New Products (2025 Launch)", "THB 82.7M invested. All early-stage GREEN on revenue trajectory. NFCV (THB 34.9M) is largest bet — lifetime projected THB 227M.", ACCENT),
    ("Forward Portfolio Value", "THB 1.06B projected lifetime revenue across 10 products vs THB 114.7M total investment. Portfolio-level ROI strong despite individual misses.", GREEN),
]

for i, (title, desc, color) in enumerate(insights):
    yi = y3 + Inches(0.38) + Inches(i * 0.39)
    # Colored bullet
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, tbl_x + Inches(0.2), yi + Inches(0.03), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # Title
    add_text(tbl_x + Inches(0.35), yi - Inches(0.03), Inches(5.3), Inches(0.17),
             title, size=8, bold=True, color=TEXT_DARK)
    add_text(tbl_x + Inches(0.35), yi + Inches(0.13), Inches(5.3), Inches(0.25),
             desc, size=7, color=TEXT_MID)


# ── Right box: Decisions Required ──
right_x = tbl_x + left_w + Inches(0.2)
right_w = Inches(6.233)
add_rect(right_x, y3, right_w, box_h, fill=WHITE, border=BORDER_GRAY)

# Red accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y3, Inches(0.06), box_h)
bar.fill.solid()
bar.fill.fore_color.rgb = RED
bar.line.fill.background()

add_text(right_x + Inches(0.25), y3 + Inches(0.1), Inches(5), Inches(0.22),
         "BOARD DECISIONS REQUIRED", size=10, bold=True, color=RED)

decisions = [
    "1.  SIC4343 — Continue, restructure, or exit?  Revenue at 3% of plan after 20 months. Investment THB 1.5M. Forward BC projects THB 26.7M lifetime if trajectory recovers. Recommendation: conduct 30-day product-market fit review before kill decision.",
    "2.  SIC7805 & Dragon — Accept extended payback or intervene?  Both products have green revenue trajectory but payback pace 2-3x above BC target. Combined investment THB 30.6M. Options: (a) accept longer payback given positive revenue, (b) restructure cost base, (c) accelerate sales pipeline.",
    "3.  GT12 (SIC279) — Approve revised ramp forecast?  Zero revenue 8 months post-launch vs THB 330K BC. Sales reports customer delays, not product issues. Need customer confirmation and revised timeline before next PLM review.",
    "4.  FOUP GP erosion — Approve pricing/COGS review?  GP at 77% vs 86% BC target on THB 64.8M revenue base. Each 1pp = ~THB 650K impact. Finance + PLM to investigate whether pricing pressure or cost structure.",
    "5.  NFCV (THB 34.9M) — Largest single investment. Confirm milestone gate criteria for D3→RFS progression.  Lifetime projected THB 227M — strong ROI if execution holds."
]

for i, d in enumerate(decisions):
    yi = y3 + Inches(0.38) + Inches(i * 0.39)
    color = RED if i < 2 else (AMBER if i < 4 else ACCENT)
    # Colored number marker
    marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x + Inches(0.2), yi + Inches(0.01), Inches(0.18), Inches(0.14))
    marker.fill.solid()
    marker.fill.fore_color.rgb = color
    marker.line.fill.background()
    marker.adjustments[0] = 0.15
    add_text(right_x + Inches(0.2), yi, Inches(0.18), Inches(0.16),
             str(i + 1), size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Decision text
    # Strip the number prefix
    txt = d.split("  ", 1)[1] if "  " in d else d
    add_text(right_x + Inches(0.45), yi - Inches(0.02), Inches(5.5), Inches(0.38),
             txt, size=7, color=TEXT_DARK)


# ── Footer ──
footer_y = Inches(7.3)
add_text(Inches(0.5), footer_y, Inches(6), Inches(0.2),
         "RAG = Trajectory vs Plan (not snapshot)  |  No Orphan Reds  |  Scoring: Revenue 40% + GP 30% + ROI 30%",
         size=6.5, color=TEXT_MID)
add_text(Inches(9), footer_y, Inches(4), Inches(0.2),
         "Source: NPD Health Dashboard, PLM Data — March 2026",
         size=6.5, color=TEXT_MID, align=PP_ALIGN.RIGHT)

# ── Save ──
output_path = "/home/user/SICT/outputs/Business-SICT/SIC-Dashboards/NPD_Health_Dashboard_Board_Slide_Mar2026.pptx"
import os
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Saved: {output_path}")
